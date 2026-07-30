import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from sqlmodel import Session, select
from collections import defaultdict
from app.models import Bill, BillItem, Product


class AnalysisService:

    def __init__(self, session: Session):
        self.session = session

    def generate(self, chat_id: int) -> str:
        bills = self.session.exec(
            select(Bill).where(Bill.chat_id == chat_id, Bill.status == "finalized")
        ).all()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # ── Title slide ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "Sales Analysis", "All finalized sales", Inches(2))

        if not bills:
            self._add_text(slide, "No sales data available.", Inches(5))
            return self._save(prs, chat_id)

        # ── Collect data ──
        total_revenue = sum(b.total_amount for b in bills)
        total_tax = sum(b.total_tax for b in bills)
        total_cgst = sum(b.total_cgst for b in bills)
        total_sgst = sum(b.total_sgst for b in bills)

        mode_totals = defaultdict(float)
        for b in bills:
            mode_totals[b.payment_mode or "unknown"] += b.total_amount

        bill_ids = [b.id for b in bills]
        item_rows = self.session.exec(select(BillItem).where(BillItem.bill_id.in_(bill_ids))).all()
        product_sales = defaultdict(lambda: {"qty": 0.0, "revenue": 0.0})
        for item in item_rows:
            product = self.session.get(Product, item.product_id)
            name = product.name.title() if product else f"Item #{item.product_id}"
            product_sales[name]["qty"] += item.quantity
            product_sales[name]["revenue"] += item.total_price
        top_items = sorted(product_sales.items(), key=lambda x: -x[1]["qty"])[:5]

        low_stock = self.session.exec(
            select(Product).where(
                Product.chat_id == chat_id, Product.stock_quantity <= Product.reorder_level
            )
        ).all()

        # ── KPI slide ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "Sales Overview", "", Inches(0.3))
        kpis = [
            ("Total Revenue", f"₹{total_revenue:,.0f}"),
            ("Bills", str(len(bills))),
            ("Tax Collected", f"₹{total_tax:,.0f}"),
            ("CGST", f"₹{total_cgst:,.0f}"),
            ("SGST", f"₹{total_sgst:,.0f}"),
        ]
        for i, (label, value) in enumerate(kpis):
            left = Inches(0.5 + i * 2.5)
            self._add_kpi(slide, left, Inches(1.5), Inches(2.3), label, value)

        # ── Pie chart ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "Revenue by Payment Mode", "", Inches(0.3))

        cd = CategoryChartData()
        cd.categories = list(mode_totals.keys())
        cd.add_series("Revenue", list(mode_totals.values()))
        cs = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(3), Inches(1.5), Inches(7), Inches(5.5), cd)
        cs.chart.has_legend = True
        cs.chart.plots[0].has_data_labels = True
        cs.chart.plots[0].data_labels.show_percentage = True
        try:
            colors = [RGBColor(0x1A, 0x56, 0xDB), RGBColor(0xE8, 0x6B, 0x1A),
                       RGBColor(0x2E, 0xCC, 0x71), RGBColor(0x9B, 0x59, 0xB6)]
            for i in range(len(mode_totals)):
                pt = cs.chart.plots[0].series[0].points[i]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        except Exception:
            pass

        # ── Quantity bar chart ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "Top Products by Quantity Sold", "", Inches(0.3))
        cd = CategoryChartData()
        cd.categories = [item[0] for item in top_items]
        cd.add_series("Qty Sold", [item[1]["qty"] for item in top_items])
        cs = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(11), Inches(5.5), cd)
        cs.chart.has_legend = False
        cs.chart.plots[0].has_data_labels = True
        try:
            cs.chart.plots[0].series[0].format.fill.solid()
            cs.chart.plots[0].series[0].format.fill.fore_color.rgb = RGBColor(0x1A, 0x56, 0xDB)
        except Exception:
            pass

        # ── Revenue bar chart ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "Revenue by Product", "", Inches(0.3))
        cd = CategoryChartData()
        cd.categories = [item[0] for item in top_items]
        cd.add_series("Revenue (₹)", [item[1]["revenue"] for item in top_items])
        cs = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(11), Inches(5.5), cd)
        cs.chart.has_legend = False
        cs.chart.plots[0].has_data_labels = True
        try:
            cs.chart.plots[0].series[0].format.fill.solid()
            cs.chart.plots[0].series[0].format.fill.fore_color.rgb = RGBColor(0x2E, 0xCC, 0x71)
        except Exception:
            pass

        # ── Stock health ──
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_title(slide, "Stock Health", "", Inches(0.3))
        if low_stock:
            rows = len(low_stock) + 1
            table_shape = slide.shapes.add_table(rows, 4, Inches(1), Inches(1.5), Inches(11), Inches(0.5 + len(low_stock) * 0.4))
            table = table_shape.table
            for col, hdr in enumerate(["Product", "Stock", "Unit", "Reorder Level"]):
                cell = table.cell(0, col)
                cell.text = hdr
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
                    p.font.size = Pt(11)
            for row, prod in enumerate(low_stock, 1):
                table.cell(row, 0).text = prod.name.title()
                table.cell(row, 1).text = str(prod.stock_quantity)
                table.cell(row, 2).text = prod.unit
                table.cell(row, 3).text = str(prod.reorder_level)
        else:
            self._add_text(slide, "All items above reorder level.", Inches(2))

        return self._save(prs, chat_id)

    def _add_title(self, slide, text, sub_text, top):
        box = slide.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(0.8))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        if sub_text:
            p2 = tf.add_paragraph()
            p2.text = sub_text
            p2.font.size = Pt(14)
            p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    def _add_kpi(self, slide, left, top, width, label, value):
        box = slide.shapes.add_textbox(left, top, width, Inches(1.2))
        tf = box.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        p2 = tf.add_paragraph()
        p2.text = value
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)

    def _add_text(self, slide, text, top):
        box = slide.shapes.add_textbox(Inches(1), top, Inches(11), Inches(1))
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x33, 0x33)

    def _save(self, prs, chat_id):
        out_dir = "generated"
        os.makedirs(out_dir, exist_ok=True)
        path = os.path.join(out_dir, f"analysis_{chat_id}.pptx")
        prs.save(path)
        return path
